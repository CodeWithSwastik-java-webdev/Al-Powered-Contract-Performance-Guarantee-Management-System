"""
Generate POWERGRID CPG Management System ER Diagram PowerPoint.
Run: python docs/generate_er_diagram_ppt.py
Output: docs/POWERGRID_CPG_ER_Diagram.pptx
"""

from __future__ import annotations

import os
from datetime import date

import matplotlib.patches as mpatches
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# POWERGRID-inspired palette
NAVY = "#003366"
BLUE = "#0066B3"
LIGHT_BLUE = "#E8F4FC"
ACCENT = "#F7941D"
WHITE = "#FFFFFF"
GRAY = "#666666"
LIGHT_GRAY = "#F5F5F5"

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PPT = os.path.join(SCRIPT_DIR, "POWERGRID_CPG_ER_Diagram.pptx")
DIAGRAM_IMG = os.path.join(SCRIPT_DIR, "_er_diagram_temp.png")


def draw_er_diagram() -> str:
    """Render ER diagram to PNG and return file path."""
    fig, ax = plt.subplots(figsize=(20, 14))
    ax.set_xlim(0, 20)
    ax.set_ylim(0, 14)
    ax.axis("off")
    fig.patch.set_facecolor(WHITE)

    def entity_box(x, y, w, h, title, fields, color=LIGHT_BLUE, border=NAVY):
        box = FancyBboxPatch(
            (x, y), w, h,
            boxstyle="round,pad=0.02,rounding_size=0.15",
            facecolor=color, edgecolor=border, linewidth=2,
        )
        ax.add_patch(box)
        ax.text(x + w / 2, y + h - 0.35, title, ha="center", va="top",
                fontsize=11, fontweight="bold", color=NAVY, family="sans-serif")
        field_text = "\n".join(fields)
        ax.text(x + 0.15, y + h - 0.7, field_text, ha="left", va="top",
                fontsize=7.5, color="#333333", family="monospace", linespacing=1.4)
        return (x + w / 2, y + h / 2)

    def arrow(x1, y1, x2, y2, label="", style="-|>", color=BLUE):
        arr = FancyArrowPatch(
            (x1, y1), (x2, y2),
            arrowstyle=style, mutation_scale=12,
            linewidth=1.5, color=color, connectionstyle="arc3,rad=0.1",
        )
        ax.add_patch(arr)
        if label:
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2
            ax.text(mx, my + 0.15, label, ha="center", va="bottom",
                    fontsize=7, color=GRAY, style="italic",
                    bbox=dict(boxstyle="round,pad=0.2", facecolor=WHITE, edgecolor="none", alpha=0.8))

    # --- Entity positions ---
    # Row 1: Users, Contractors
    users_c = entity_box(0.3, 10.5, 3.2, 2.8, "USERS", [
        "PK  id",
        "UK  firebase_uid",
        "    name, email",
        "    role (enum)",
        "    department, phone",
        "    is_active",
        "    created_at, updated_at",
    ])

    contractors_c = entity_box(7.5, 10.5, 3.5, 2.8, "CONTRACTORS", [
        "PK  id",
        "UK  vendor_code",
        "    name, gst_number",
        "    pan_number",
        "    contact_person",
        "    email, phone, rating",
        "    is_active, is_blacklisted",
    ])

    # Row 2: Contracts (center)
    contracts_c = entity_box(5.5, 7.0, 4.0, 2.6, "CONTRACTS", [
        "PK  id",
        "UK  contract_number",
        "    project_name, description",
        "    contract_value, currency",
        "    award_date, completion_date",
        "    status (enum), zone",
        "FK  contractor_id",
        "FK  created_by_id, updated_by_id",
    ])

    # Row 3: CPGs (center)
    cpgs_c = entity_box(5.5, 3.8, 4.0, 2.8, "CPGS", [
        "PK  id",
        "UK  bg_number",
        "FK  contract_id",
        "    bg_type, bank_name",
        "    amount, issue_date",
        "    expiry_date, claim_period_end",
        "    status (enum)",
        "FK  renewed_from_id (self)",
    ])

    # Row 3 sides: Documents, Risk, ML
    docs_c = entity_box(0.3, 3.5, 3.2, 2.6, "DOCUMENTS", [
        "PK  id",
        "FK  cpg_id",
        "    file_name, file_url",
        "    cloudinary_public_id",
        "    document_type (enum)",
        "FK  uploaded_by",
        "    version, is_active",
    ])

    risk_c = entity_box(11.0, 3.5, 3.5, 2.6, "RISK_ASSESSMENTS", [
        "PK  id",
        "FK  cpg_id",
        "    health_score",
        "    risk_level (enum)",
        "    anomaly_detected",
        "    anomaly_reason",
        "    factors (JSON)",
    ])

    ml_c = entity_box(15.2, 3.5, 3.5, 2.6, "ML_PREDICTIONS", [
        "PK  id",
        "FK  cpg_id",
        "    risk_probability",
        "    delay_probability",
        "    model_version",
        "    confidence_score",
        "    features (JSON)",
    ])

    # Row 4: Audit, Notifications
    audit_c = entity_box(0.3, 0.3, 3.5, 2.4, "AUDIT_LOGS", [
        "PK  id",
        "FK  user_id (nullable)",
        "    entity_type (enum)",
        "    entity_id",
        "    action (enum)",
        "    old_value, new_value (JSON)",
        "    ip_address, user_agent",
    ])

    notif_c = entity_box(7.5, 0.3, 3.5, 2.4, "NOTIFICATIONS", [
        "PK  id",
        "FK  user_id",
        "    type (enum)",
        "    title, message",
        "    is_read, read_at",
        "    related_entity_type",
        "    related_entity_id",
    ])

    # --- Relationships ---
    arrow(contractors_c[0], contractors_c[1] - 1.2, contracts_c[0] + 0.5, contracts_c[1] + 1.0, "1 : N")
    arrow(contracts_c[0], contracts_c[1] - 1.1, cpgs_c[0], cpgs_c[1] + 1.2, "1 : N")
    arrow(cpgs_c[0] - 1.5, cpgs_c[1], docs_c[0] + 1.2, docs_c[1], "1 : N")
    arrow(cpgs_c[0] + 1.5, cpgs_c[1], risk_c[0] - 1.2, risk_c[1], "1 : N")
    arrow(cpgs_c[0] + 2.0, cpgs_c[1] + 0.3, ml_c[0] - 1.5, ml_c[1], "1 : N")

    # User relations
    arrow(users_c[0] + 0.8, users_c[1] - 1.2, contracts_c[0] - 1.2, contracts_c[1] + 1.0, "creates/updates")
    arrow(users_c[0] + 0.5, users_c[1] - 1.4, docs_c[0] + 0.5, docs_c[1] + 1.1, "uploads")
    arrow(users_c[0], users_c[1] - 1.4, audit_c[0], audit_c[1] + 1.0, "1 : N")
    arrow(users_c[0] + 1.5, users_c[1] - 1.4, notif_c[0], notif_c[1] + 1.0, "1 : N")

    # Self-referential CPG renewal
    ax.annotate(
        "", xy=(cpgs_c[0] + 2.3, cpgs_c[1] + 1.1), xytext=(cpgs_c[0] + 2.3, cpgs_c[1] + 0.2),
        arrowprops=dict(arrowstyle="-|>", color=ACCENT, lw=1.5,
                        connectionstyle="arc3,rad=-0.4"),
    )
    ax.text(cpgs_c[0] + 2.8, cpgs_c[1] + 0.7, "renewal\n(self-ref)", fontsize=7,
            color=ACCENT, ha="left", style="italic")

    # Title
    ax.text(10, 13.5, "POWERGRID — CPG Management System", ha="center", va="center",
            fontsize=18, fontweight="bold", color=NAVY)
    ax.text(10, 13.0, "Entity-Relationship Diagram  |  PostgreSQL + Prisma ORM",
            ha="center", va="center", fontsize=11, color=GRAY)

    # Legend
    legend_y = 0.5
    ax.add_patch(FancyBboxPatch((14.5, legend_y), 5.0, 2.2,
                                boxstyle="round,pad=0.1", facecolor=LIGHT_GRAY, edgecolor=GRAY, linewidth=1))
    ax.text(17.0, legend_y + 1.9, "Legend", ha="center", fontsize=9, fontweight="bold", color=NAVY)
    ax.text(14.7, legend_y + 1.45, "PK = Primary Key", fontsize=8, color="#333")
    ax.text(14.7, legend_y + 1.15, "FK = Foreign Key", fontsize=8, color="#333")
    ax.text(14.7, legend_y + 0.85, "UK = Unique Key", fontsize=8, color="#333")
    ax.text(14.7, legend_y + 0.55, "Restrict = parent delete blocked", fontsize=8, color="#333")
    ax.text(14.7, legend_y + 0.25, "Cascade = child records deleted", fontsize=8, color="#333")

    plt.tight_layout(pad=0.5)
    plt.savefig(DIAGRAM_IMG, dpi=180, bbox_inches="tight", facecolor=WHITE)
    plt.close()
    return DIAGRAM_IMG


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, RGBColor(0, 51, 102))

    # Accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(5.8), Inches(10), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(247, 148, 29)
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "POWERGRID"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.text = "Contract Performance Guarantee (CPG)"
    p2.font.size = Pt(28)
    p2.font.color.rgb = RGBColor(232, 244, 252)
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "Management System — Database ER Diagram"
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(200, 220, 240)
    p3.space_before = Pt(4)

    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8), Inches(0.5))
    fp = footer.text_frame.paragraphs[0]
    fp.text = f"PostgreSQL  •  Prisma ORM  •  Railway  •  {date.today().strftime('%B %Y')}"
    fp.font.size = Pt(12)
    fp.font.color.rgb = RGBColor(180, 200, 220)


def add_er_diagram_slide(prs: Presentation, img_path: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))

    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.5))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Complete Entity-Relationship Diagram"
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0, 51, 102)

    slide.shapes.add_picture(img_path, Inches(0.15), Inches(0.55), width=Inches(9.7))


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))

    # Header bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
        sp = st.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(102, 102, 102)
        top = Inches(1.5)
    else:
        top = Inches(1.2)

    body = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)
        p.level = 0


def add_relationship_table_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tp = tb.text_frame.paragraphs[0]
    tp.text = "Relationships & Cascade Rules"
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(255, 255, 255)

    rows, cols = 11, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.1), Inches(9.4), Inches(5.8))
    table = table_shape.table

    headers = ["Parent", "Child", "Cardinality", "onDelete"]
    data = [
        ["Contractor", "Contract", "1 : N", "Restrict"],
        ["Contract", "Cpg", "1 : N", "Restrict"],
        ["Cpg", "Document", "1 : N", "Cascade"],
        ["Cpg", "RiskAssessment", "1 : N", "Cascade"],
        ["Cpg", "MlPrediction", "1 : N", "Cascade"],
        ["Cpg", "Cpg (renewal)", "1 : N", "SetNull"],
        ["User", "Document", "1 : N", "Restrict"],
        ["User", "AuditLog", "1 : N", "SetNull"],
        ["User", "Notification", "1 : N", "Cascade"],
        ["User", "Contract", "1 : N", "SetNull"],
    ]

    col_widths = [Inches(1.8), Inches(2.2), Inches(1.5), Inches(1.5)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 179)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r, row in enumerate(data, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(232, 244, 252)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_enums_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(255, 255, 255))

    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0, 51, 102)
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tp = tb.text_frame.paragraphs[0]
    tp.text = "Prisma Enums"
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(255, 255, 255)

    enums = [
        ("UserRole", "ADMIN, ENGINEER, FINANCE, VIEWER"),
        ("ContractStatus", "DRAFT, ACTIVE, COMPLETED, TERMINATED, SUSPENDED"),
        ("CpgStatus", "ACTIVE, EXPIRED, CLAIMED, INVOKED, RELEASED, RENEWED, CANCELLED"),
        ("BankGuaranteeType", "PERFORMANCE_BANK_GUARANTEE, ADVANCE_BANK_GUARANTEE, SECURITY_DEPOSIT, OTHER"),
        ("DocumentType", "BANK_GUARANTEE, CONTRACT_COPY, EXTENSION_LETTER, CLAIM_NOTICE, ..."),
        ("RiskLevel", "LOW, MEDIUM, HIGH, CRITICAL"),
        ("AuditAction", "CREATE, UPDATE, DELETE, STATUS_CHANGE, UPLOAD, INVOCATION, ..."),
        ("NotificationType", "EXPIRY_ALERT, ANOMALY_DETECTED, STATUS_CHANGE, RISK_ESCALATION, ..."),
    ]

    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    tf = left.text_frame
    tf.word_wrap = True
    for i, (name, values) in enumerate(enums):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = name
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = f"   {values}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(80, 80, 80)
        p2.space_after = Pt(10)


def build_presentation() -> str:
    print("Drawing ER diagram...")
    img_path = draw_er_diagram()

    print("Building PowerPoint...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_er_diagram_slide(prs, img_path)
    add_bullet_slide(
        prs,
        "Architecture Overview",
        [
            "9 core tables modelling the full CPG lifecycle for POWERGRID",
            "PostgreSQL on Railway with Prisma ORM for type-safe data access",
            "Firebase Authentication — users linked via firebase_uid",
            "Cloudinary for document storage — metadata tracked in documents table",
            "Append-only AI tables: risk_assessments & ml_predictions (time-series per CPG)",
            "Immutable audit_logs for compliance and forensic review",
            "Polymorphic references in audit_logs & notifications (entity_type + entity_id)",
        ],
        "Tech Stack: Node.js • Express • TypeScript • Prisma • PostgreSQL",
    )
    add_bullet_slide(
        prs,
        "Core Domain Hierarchy",
        [
            "Contractor  →  owns  →  Contract  →  governs  →  Cpg",
            "Cpg  →  has many  →  Documents (Cloudinary files)",
            "Cpg  →  has many  →  RiskAssessments (Health Score Engine)",
            "Cpg  →  has many  →  MlPredictions (ML Risk Prediction)",
            "Cpg  →  self-renews via  →  renewed_from_id (renewal chain)",
            "User  →  uploads  →  Documents",
            "User  →  receives  →  Notifications (expiry alerts, anomalies)",
            "User  →  generates  →  AuditLogs (all CRUD & status changes)",
        ],
    )
    add_relationship_table_slide(prs)
    add_enums_slide(prs)
    add_bullet_slide(
        prs,
        "Key Indexes (Performance)",
        [
            "cpgs (status, expiry_date) — expiry dashboard & alert cron jobs",
            "cpgs (claim_period_end) — claim window monitoring",
            "notifications (user_id, is_read) — unread badge & inbox queries",
            "audit_logs (entity_type, entity_id) — entity history lookup",
            "risk_assessments (cpg_id, created_at) — latest health score per CPG",
            "ml_predictions (cpg_id, created_at) — latest ML inference per CPG",
            "contracts (contractor_id) — all contracts for a vendor",
            "users (firebase_uid) — Firebase auth sync on every request",
        ],
    )
    add_bullet_slide(
        prs,
        "Phase 2 — Recommended Extensions",
        [
            "cpg_status_history — immutable status transition log with actor & reason",
            "cpg_claims — formal BG invocation workflow (claim amount, bank response)",
            "cpg_renewal_requests — Engineer → Finance → Admin approval workflow",
            "expiry_alert_rules — configurable 90/60/30-day thresholds per zone",
            "banks — master bank table with IFSC validation",
            "user_zone_assignments — scoped RBAC (Engineer sees only NR zone)",
            "ai_assistant_sessions — chat history for AI Assistant feature",
        ],
        "Not yet in schema — add when features are implemented",
    )

    prs.save(OUTPUT_PPT)

    if os.path.exists(img_path):
        os.remove(img_path)

    return OUTPUT_PPT


if __name__ == "__main__":
    out = build_presentation()
    print(f"Done: {out}")
